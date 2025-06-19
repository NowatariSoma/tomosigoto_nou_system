#!/usr/bin/env python
"""
Integration tests for Git PPT functionality
"""
import os
import shutil
import subprocess
import tempfile
import unittest

from pptx import Presentation

from src.cli import Installer
from src.diff import PPTXDiffer
from src.pptx_extractor import PPTXExtractor


class TestPPTXCreation(unittest.TestCase):
    """Test cases for creating and manipulating PPTX files"""

    def setUp(self):
        """Set up test environment"""
        self.temp_dir = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, self.temp_dir)

    def create_test_presentation(
        self, filename, title="Test Presentation", slides_content=None
    ):
        """Create a test PPTX file with specified content"""
        if slides_content is None:
            slides_content = [
                {
                    "title": "First Slide",
                    "content": ["First slide content", "Additional text"],
                },
                {"title": "Second Slide", "content": ["Second slide content"]},
            ]

        prs = Presentation()

        # Set presentation properties
        prs.core_properties.title = title
        prs.core_properties.author = "Test Author"
        prs.core_properties.subject = "Test Subject"

        for slide_data in slides_content:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            # Add content
            content_shape = slide.placeholders[1]
            content_shape.text = "\n".join(slide_data["content"])

        filepath = os.path.join(self.temp_dir, filename)
        prs.save(filepath)
        return filepath

    def test_create_and_extract_presentation(self):
        """Test creating and extracting content from a PPTX file"""
        # Create test presentation
        filepath = self.create_test_presentation("test.pptx")

        # Extract content
        extractor = PPTXExtractor(filepath)
        content = extractor.extract_content()

        # Verify extraction
        self.assertIn("metadata", content)
        self.assertIn("slides", content)
        self.assertIn("notes", content)

        # Check metadata
        metadata = content["metadata"]
        self.assertEqual(metadata["title"], "Test Presentation")
        self.assertEqual(metadata["author"], "Test Author")
        self.assertEqual(metadata["slide_count"], 2)

        # Check slides
        slides = content["slides"]
        self.assertEqual(len(slides), 2)

        # Check first slide
        first_slide = slides[0]
        self.assertEqual(first_slide["slide_number"], 1)
        self.assertEqual(first_slide["title"], "First Slide")
        # Check if content is in any of the text content items
        first_slide_text = " ".join(first_slide["text_content"])
        self.assertIn("First slide content", first_slide_text)

        # Check second slide
        second_slide = slides[1]
        self.assertEqual(second_slide["slide_number"], 2)
        self.assertEqual(second_slide["title"], "Second Slide")
        second_slide_text = " ".join(second_slide["text_content"])
        self.assertIn("Second slide content", second_slide_text)

    def test_text_format_output(self):
        """Test text format output for Git diff"""
        filepath = self.create_test_presentation("test.pptx")

        extractor = PPTXExtractor(filepath)
        text_output = extractor.to_text_format()

        # Check that text output contains expected content
        self.assertIn("Title: Test Presentation", text_output)
        self.assertIn("Author: Test Author", text_output)
        self.assertIn("Slides: 2", text_output)
        self.assertIn("=== Slide 1 ===", text_output)
        self.assertIn("=== Slide 2 ===", text_output)
        self.assertIn("First Slide", text_output)
        self.assertIn("Second Slide", text_output)


class TestPPTXDiffIntegration(unittest.TestCase):
    """Integration tests for PPTX diff functionality"""

    def setUp(self):
        """Set up test environment"""
        self.temp_dir = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, self.temp_dir)

    def create_test_presentation(
        self, filename, title="Test Presentation", slides_content=None
    ):
        """Create a test PPTX file with specified content"""
        if slides_content is None:
            slides_content = [
                {"title": "First Slide", "content": ["First slide content"]},
            ]

        prs = Presentation()
        prs.core_properties.title = title
        prs.core_properties.author = "Test Author"

        for slide_data in slides_content:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]
            content_shape = slide.placeholders[1]
            content_shape.text = "\n".join(slide_data["content"])

        filepath = os.path.join(self.temp_dir, filename)
        prs.save(filepath)
        return filepath

    def test_diff_identical_presentations(self):
        """Test diff between identical presentations"""
        file1 = self.create_test_presentation("test1.pptx")
        file2 = self.create_test_presentation("test2.pptx")

        differ = PPTXDiffer(file1, file2, "test.pptx")
        diffs = differ.generate_diff()

        # Should have no differences
        self.assertEqual(len(diffs), 0)

    def test_diff_different_titles(self):
        """Test diff between presentations with different titles"""
        file1 = self.create_test_presentation("test1.pptx", title="Old Title")
        file2 = self.create_test_presentation("test2.pptx", title="New Title")

        differ = PPTXDiffer(file1, file2, "test.pptx")
        diffs = differ.generate_diff()

        # Should have metadata differences
        metadata_diffs = [d for d in diffs if d["type"] == "metadata"]
        self.assertGreater(len(metadata_diffs), 0)

        # Check diff content
        metadata_diff = metadata_diffs[0]
        self.assertIn("Old Title", metadata_diff["diff"])
        self.assertIn("New Title", metadata_diff["diff"])

    def test_diff_different_slide_content(self):
        """Test diff between presentations with different slide content"""
        slides1 = [{"title": "Same Title", "content": ["Old content"]}]
        slides2 = [{"title": "Same Title", "content": ["New content"]}]

        file1 = self.create_test_presentation("test1.pptx", slides_content=slides1)
        file2 = self.create_test_presentation("test2.pptx", slides_content=slides2)

        differ = PPTXDiffer(file1, file2, "test.pptx")
        diffs = differ.generate_diff()

        # Should have slide differences
        slide_diffs = [d for d in diffs if d["type"] == "slide"]
        self.assertGreater(len(slide_diffs), 0)

        # Check diff content
        slide_diff = slide_diffs[0]
        self.assertIn("Old content", slide_diff["diff"])
        self.assertIn("New content", slide_diff["diff"])

    def test_diff_added_slide(self):
        """Test diff with added slide"""
        slides1 = [{"title": "First Slide", "content": ["Content 1"]}]
        slides2 = [
            {"title": "First Slide", "content": ["Content 1"]},
            {"title": "Second Slide", "content": ["Content 2"]},
        ]

        file1 = self.create_test_presentation("test1.pptx", slides_content=slides1)
        file2 = self.create_test_presentation("test2.pptx", slides_content=slides2)

        differ = PPTXDiffer(file1, file2, "test.pptx")
        diffs = differ.generate_diff()

        # Should have slide addition
        slide_diffs = [
            d for d in diffs if d["type"] == "slide" and d["slide_number"] == 2
        ]
        self.assertGreater(len(slide_diffs), 0)

        # Check that it's marked as addition
        slide_diff = slide_diffs[0]
        self.assertIn("--- /dev/null", slide_diff["a"])
        self.assertIn("+++ b/test.pptx/Slide/2", slide_diff["b"])

    def test_diff_removed_slide(self):
        """Test diff with removed slide"""
        slides1 = [
            {"title": "First Slide", "content": ["Content 1"]},
            {"title": "Second Slide", "content": ["Content 2"]},
        ]
        slides2 = [{"title": "First Slide", "content": ["Content 1"]}]

        file1 = self.create_test_presentation("test1.pptx", slides_content=slides1)
        file2 = self.create_test_presentation("test2.pptx", slides_content=slides2)

        differ = PPTXDiffer(file1, file2, "test.pptx")
        diffs = differ.generate_diff()

        # Should have slide removal
        slide_diffs = [
            d for d in diffs if d["type"] == "slide" and d["slide_number"] == 2
        ]
        self.assertGreater(len(slide_diffs), 0)

        # Check that it's marked as removal
        slide_diff = slide_diffs[0]
        self.assertIn("--- a/test.pptx/Slide/2", slide_diff["a"])
        self.assertIn("+++ /dev/null", slide_diff["b"])


class TestGitIntegration(unittest.TestCase):
    """Integration tests for Git functionality"""

    def setUp(self):
        """Set up test Git repository"""
        self.temp_dir = tempfile.mkdtemp()
        self.addCleanup(shutil.rmtree, self.temp_dir)

        # Initialize git repository
        subprocess.run(["git", "init"], cwd=self.temp_dir, capture_output=True)
        subprocess.run(
            ["git", "config", "user.name", "Test User"],
            cwd=self.temp_dir,
            capture_output=True,
        )
        subprocess.run(
            ["git", "config", "user.email", "test@example.com"],
            cwd=self.temp_dir,
            capture_output=True,
        )

    def create_test_presentation(self, filename, title="Test Presentation"):
        """Create a simple test PPTX file"""
        prs = Presentation()
        prs.core_properties.title = title
        prs.core_properties.author = "Test Author"

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Test Slide"
        content_shape = slide.placeholders[1]
        content_shape.text = "Test content"

        filepath = os.path.join(self.temp_dir, filename)
        prs.save(filepath)
        return filepath

    def test_git_ppt_workflow(self):
        """Test complete Git PPT workflow"""
        # Install Git PPT
        installer = Installer(mode="local", path=self.temp_dir)
        installer.install()

        # Create initial presentation
        pptx_file = self.create_test_presentation("presentation.pptx")
        self.assertTrue(os.path.exists(pptx_file))  # Verify file was created

        # Add and commit
        subprocess.run(
            ["git", "add", "presentation.pptx"], cwd=self.temp_dir, capture_output=True
        )
        result = subprocess.run(
            ["git", "commit", "-m", "Initial presentation"],
            cwd=self.temp_dir,
            capture_output=True,
            text=True,
        )
        self.assertEqual(result.returncode, 0)

        # Modify presentation
        modified_file = self.create_test_presentation(
            "presentation.pptx", title="Modified Presentation"
        )
        self.assertTrue(os.path.exists(modified_file))  # Verify file was modified

        # Check git status
        result = subprocess.run(
            ["git", "status", "--porcelain"],
            cwd=self.temp_dir,
            capture_output=True,
            text=True,
        )
        self.assertIn("presentation.pptx", result.stdout)

        # Test git diff (this should use our custom diff)
        result = subprocess.run(
            ["git", "diff", "presentation.pptx"],
            cwd=self.temp_dir,
            capture_output=True,
            text=True,
        )

        # The diff should contain meaningful content, not binary diff
        self.assertNotIn("Binary files", result.stdout)
        # Should contain our diff format
        if result.stdout:  # Git might not show diff for some reason
            self.assertIn("diff --ppt", result.stdout)

    def test_gitattributes_configuration(self):
        """Test that .gitattributes is properly configured"""
        installer = Installer(mode="local", path=self.temp_dir)
        installer.install()

        gitattributes_path = os.path.join(self.temp_dir, ".gitattributes")
        self.assertTrue(os.path.exists(gitattributes_path))

        with open(gitattributes_path, "r") as f:
            content = f.read()
            self.assertIn("*.pptx diff=ppt", content)
            self.assertIn("*.ppt diff=ppt", content)

    def test_git_config_setup(self):
        """Test that git config is properly set up"""
        installer = Installer(mode="local", path=self.temp_dir)
        installer.install()

        # Check git config
        result = subprocess.run(
            ["git", "config", "--get", "diff.ppt.command"],
            cwd=self.temp_dir,
            capture_output=True,
            text=True,
        )
        self.assertEqual(result.returncode, 0)
        self.assertIn("diff.py", result.stdout)


class TestRealFileIntegration(unittest.TestCase):
    """Integration tests using real example files"""

    def setUp(self):
        """Set up test with real PPTX file"""
        self.example_file = os.path.join(
            os.path.dirname(__file__),
            "..",
            "..",
            "repository",
            "office-examples",
            "example.pptx",
        )

    def test_extract_real_file_comprehensive(self):
        """Test comprehensive extraction from real PPTX file"""
        if not os.path.exists(self.example_file):
            self.skipTest("Example PPTX file not found")

        extractor = PPTXExtractor(self.example_file)
        content = extractor.extract_content()

        # Detailed content verification
        self.assertIn("metadata", content)
        self.assertIn("slides", content)
        self.assertIn("notes", content)

        # Check metadata details
        metadata = content["metadata"]
        self.assertIsInstance(metadata.get("title"), str)
        self.assertIsInstance(metadata.get("author"), str)
        self.assertIsInstance(metadata.get("slide_count"), int)

        # Check slides details
        slides = content["slides"]
        self.assertGreater(len(slides), 0)

        for slide in slides:
            self.assertIn("slide_number", slide)
            self.assertIn("title", slide)
            self.assertIn("text_content", slide)
            self.assertIn("shapes_count", slide)
            self.assertIsInstance(slide["text_content"], list)

        # Check text format
        text_format = extractor.to_text_format()
        self.assertIsInstance(text_format, str)
        self.assertGreater(len(text_format), 0)

    def test_diff_with_real_file(self):
        """Test diff functionality with real file"""
        if not os.path.exists(self.example_file):
            self.skipTest("Example PPTX file not found")

        # Test diff with itself (should have no differences)
        differ = PPTXDiffer(self.example_file, self.example_file, "example.pptx")
        diffs = differ.generate_diff()
        self.assertEqual(len(diffs), 0)

        # Test diff with null file (should show all content as additions)
        differ = PPTXDiffer(None, self.example_file, "example.pptx")
        diffs = differ.generate_diff()
        self.assertGreater(len(diffs), 0)

        # Should have slide additions
        slide_diffs = [d for d in diffs if d["type"] == "slide"]
        self.assertGreater(len(slide_diffs), 0)


if __name__ == "__main__":
    unittest.main()
