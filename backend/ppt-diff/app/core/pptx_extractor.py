#!/usr/bin/env python
"""
PowerPoint content extractor using python-pptx
"""
import json
import os
import sys
import tempfile
from typing import Dict, List, Optional, Tuple

from pptx import Presentation


class PPTXExtractor:
    """Extract content from PowerPoint files for Git diff purposes"""

    def __init__(self, file_path: str):
        self.file_path = file_path
        self.presentation = None

    def extract_content(self) -> Dict:
        """Extract all content from the PowerPoint file"""
        if not os.path.exists(self.file_path):
            return {}

        try:
            # Use python-pptx for structured extraction
            self.presentation = Presentation(self.file_path)
            content = {
                "metadata": self._extract_metadata(),
                "slides": self._extract_slides(),
                "notes": self._extract_notes(),
            }
            return content
        except Exception as e:
            print(
                f"Error extracting content from {self.file_path}: {e}", file=sys.stderr
            )
            return {}

    def _extract_metadata(self) -> Dict:
        """Extract presentation metadata"""
        if not self.presentation:
            return {}

        core_props = self.presentation.core_properties
        return {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "slide_count": len(self.presentation.slides),
        }

    def _extract_slides(self) -> List[Dict]:
        """Extract content from each slide"""
        if not self.presentation:
            return []

        slides = []
        for i, slide in enumerate(self.presentation.slides, 1):
            slide_content = {
                "slide_number": i,
                "title": self._get_slide_title(slide),
                "text_content": self._get_slide_text(slide),
                "shapes_count": len(slide.shapes),
            }
            slides.append(slide_content)
        return slides

    def _extract_notes(self) -> List[Dict]:
        """Extract speaker notes from slides"""
        if not self.presentation:
            return []

        notes = []
        for i, slide in enumerate(self.presentation.slides, 1):
            if slide.has_notes_slide:
                notes_text = self._get_notes_text(slide.notes_slide)
                if notes_text.strip():
                    notes.append({"slide_number": i, "notes": notes_text})
        return notes

    def _get_slide_title(self, slide) -> str:
        """Extract title from slide"""
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Assume first text shape is title
                title = shape.text.strip()
                break
        return title

    def _get_slide_text(self, slide) -> List[str]:
        """Extract all text content from slide"""
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        return text_content

    def _get_notes_text(self, notes_slide) -> str:
        """Extract text from notes slide"""
        notes_text = ""
        for shape in notes_slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                notes_text += shape.text.strip() + "\n"
        return notes_text.strip()

    def to_text_format(self) -> str:
        """Convert extracted content to text format for diff"""
        content = self.extract_content()
        if not content:
            return ""

        lines = []

        # Add metadata
        metadata = content.get("metadata", {})
        if metadata.get("title"):
            lines.append(f"Title: {metadata['title']}")
        if metadata.get("author"):
            lines.append(f"Author: {metadata['author']}")
        if metadata.get("slide_count"):
            lines.append(f"Slides: {metadata['slide_count']}")

        if lines:
            lines.append("")  # Empty line separator

        # Add slide content
        slides = content.get("slides", [])
        for slide in slides:
            lines.append(f"=== Slide {slide['slide_number']} ===")
            if slide.get("title"):
                lines.append(f"Title: {slide['title']}")

            text_content = slide.get("text_content", [])
            for text in text_content:
                if text.strip():
                    lines.append(text.strip())

            lines.append("")  # Empty line separator

        # Add notes
        notes = content.get("notes", [])
        if notes:
            lines.append("=== Speaker Notes ===")
            for note in notes:
                lines.append(f"Slide {note['slide_number']} Notes:")
                lines.append(note["notes"])
                lines.append("")

        return "\n".join(lines)


def extract_pptx_content(file_path: str) -> str:
    """Main function to extract PPTX content as text"""
    extractor = PPTXExtractor(file_path)
    return extractor.to_text_format()


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python pptx_extractor.py <pptx_file>")
        sys.exit(1)

    file_path = sys.argv[1]
    content = extract_pptx_content(file_path)
    print(content)
